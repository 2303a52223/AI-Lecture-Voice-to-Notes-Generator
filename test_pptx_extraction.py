#!/usr/bin/env python3
"""
Test PPTX extraction directly
"""
import sys
from pathlib import Path

print("=" * 70)
print("PPTX EXTRACTION TEST")
print("=" * 70)

# Check file
pptx_file = Path("data/uploads/20260502_105844_Milestone2_TextImageGenAI.pptx")

if not pptx_file.exists():
    print(f"✗ File not found: {pptx_file}")
    sys.exit(1)

print(f"\n✓ File found: {pptx_file.name}")
print(f"  Size: {pptx_file.stat().st_size} bytes")

# Test python-pptx
print("\nTesting python-pptx library...")
try:
    from pptx import Presentation
    
    prs = Presentation(str(pptx_file))
    print(f"✓ File opened successfully")
    print(f"  Total slides: {len(prs.slides)}")
    
    all_text = ""
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n  Slide {slide_idx + 1}:")
        slide_text = ""
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                print(f"    Shape {shape_idx}: {text[:80]}...")
                slide_text += text + "\n"
        
        if not slide_text:
            print(f"    (no text found)")
        
        all_text += slide_text
    
    print(f"\n✓ Total extracted text: {len(all_text)} characters")
    if all_text.strip():
        print(f"  Preview: {all_text[:300]}...")
    else:
        print("  ✗ No text was extracted - PPTX may be image-based")
        
except ImportError:
    print("✗ python-pptx not installed")
except Exception as e:
    print(f"✗ Error: {type(e).__name__}: {e}")
