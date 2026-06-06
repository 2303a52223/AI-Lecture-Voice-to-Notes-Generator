"""
Document extractor stub

Provides simple extraction for PDF, PPTX, and DOCX files and a routing helper.
This is a minimal, local-first implementation that attempts to use `fitz` (PyMuPDF),
`python-pptx`, and `python-docx` when available. OCR is attempted via `pytesseract`
if requested and pillow is available.

Return shape:
{
    'text': str,
    'segments': [{'id': int, 'page': int, 'text': str}],
    'metadata': {'pages': int, 'slides': int, ...}
}
"""
from pathlib import Path
import json
from utils.retry import retry_call
from utils.error_handler import report_error


def extract_from_pdf(file_path, ocr=False):
    file_path = Path(file_path)
    segments = []
    fitz = None
    
    # Try to import PyMuPDF
    try:
        import fitz  # PyMuPDF
    except ImportError:
        try:
            import pymupdf as fitz
        except ImportError:
            pass
    
    # Try alternative: pypdf
    pypdf_available = False
    try:
        import pypdf
        pypdf_available = True
    except ImportError:
        pass
    
    # Try pdfplumber as last resort
    pdfplumber_available = False
    try:
        import pdfplumber
        pdfplumber_available = True
    except ImportError:
        pass

    try:
        # Primary method: PyMuPDF (fitz)
        if fitz:
            try:
                try:
                    doc = retry_call(lambda: fitz.open(str(file_path)), tries=3, delay=0.5, backoff=2.0)
                except Exception as e:
                    report_error(e, f"Failed to open PDF with PyMuPDF: {file_path}", user_facing=False)
                    raise

                for i, page in enumerate(doc):
                    try:
                        text = page.get_text("text").strip()
                    except Exception:
                        text = ""
                    
                    # If no text extracted, try OCR if enabled
                    if not text and ocr:
                        try:
                            from PIL import Image
                            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                            import pytesseract
                            text = pytesseract.image_to_string(img)
                        except Exception as e:
                            text = f'[OCR failed: {str(e)[:50]}]'
                    
                    if text:  # Only add if we got text
                        segments.append({'id': i, 'page': i + 1, 'text': text})
                
                full_text = "\n\n".join(s['text'] for s in segments if s['text'])
                if full_text:
                    return {'text': full_text, 'segments': segments, 'metadata': {'pages': len(doc), 'method': 'pymupdf'}}
            except Exception as e:
                # Fall through to next method
                pass
        
        # Fallback: pdfplumber (better for scanned PDFs)
        if pdfplumber_available and not segments:
            try:
                import pdfplumber
                try:
                    pdf = retry_call(lambda: pdfplumber.open(str(file_path)), tries=2, delay=0.5)
                except Exception as e:
                    report_error(e, f"Failed to open PDF with pdfplumber: {file_path}", user_facing=False)
                    raise

                with pdf:
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text() or ""
                        text = text.strip()
                        if text:
                            segments.append({'id': i, 'page': i + 1, 'text': text})

                full_text = "\n\n".join(s['text'] for s in segments if s['text'])
                if full_text:
                    return {'text': full_text, 'segments': segments, 'metadata': {'pages': len(pdf.pages), 'method': 'pdfplumber'}}
            except Exception as e:
                # continue to other fallbacks
                pass
        
        # Second fallback: pypdf library
        if pypdf_available and not segments:
            try:
                from pypdf import PdfReader
                try:
                    reader = retry_call(lambda: PdfReader(str(file_path)), tries=2, delay=0.5)
                except Exception as e:
                    report_error(e, f"Failed to read PDF with pypdf: {file_path}", user_facing=False)
                    raise

                page_count = len(reader.pages)

                for i, page in enumerate(reader.pages):
                    try:
                        text = page.extract_text().strip() if page.extract_text() else ""
                    except Exception:
                        text = ""
                    if text:
                        segments.append({'id': i, 'page': i + 1, 'text': text})

                full_text = "\n\n".join(s['text'] for s in segments if s['text'])
                if full_text:
                    return {'text': full_text, 'segments': segments, 'metadata': {'pages': page_count, 'method': 'pypdf'}}
            except Exception as e:
                pass
        
        # If still no text, return error
        if not segments:
            methods_tried = []
            if fitz: methods_tried.append('pymupdf')
            if pdfplumber_available: methods_tried.append('pdfplumber')
            if pypdf_available: methods_tried.append('pypdf')
            
            error_msg = f"No text extracted. Methods tried: {', '.join(methods_tried) if methods_tried else 'none'}. "
            error_msg += "PDF may be image-based. Enable OCR or convert to a text-based PDF."
            
            return {'text': '', 'segments': [], 'metadata': {'error': error_msg, 'methods_tried': methods_tried}}
        
        return {'text': '', 'segments': [], 'metadata': {'error': 'All extraction methods failed'}}

    except Exception as e:
        report_error(e, f"PDF extraction exception for {file_path}", user_facing=False)
        return {'text': '', 'segments': [], 'metadata': {'error': f'PDF extraction exception: {str(e)}'}}


def extract_from_pptx(file_path):
    file_path = Path(file_path)
    segments = []
    try:
        from pptx import Presentation
    except Exception:
        return {'text': '', 'segments': [], 'metadata': {'slides': 0, 'note': 'python-pptx not available'}}

    try:
        prs = Presentation(str(file_path))
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    texts.append(shape.text)
            slide_text = '\n'.join(t for t in texts if t)
            segments.append({'id': i, 'slide': i + 1, 'text': slide_text})
        full_text = "\n\n".join(s['text'] for s in segments if s['text'])
        return {'text': full_text, 'segments': segments, 'metadata': {'slides': len(prs.slides)}}
    except Exception as e:
        return {'text': '', 'segments': [], 'metadata': {'error': str(e)}}


def extract_from_docx(file_path):
    file_path = Path(file_path)
    try:
        import docx
    except Exception:
        return {'text': '', 'segments': [], 'metadata': {'note': 'python-docx not available'}}

    try:
        doc = docx.Document(str(file_path))
        paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]

        # DOCX has no reliable page boundaries in python-docx.
        # Group paragraphs into "logical pages" so we can process one unit at a time.
        page_size = 8
        segments = []
        for i in range(0, len(paras), page_size):
            chunk = paras[i:i + page_size]
            if not chunk:
                continue
            page_num = (i // page_size) + 1
            segments.append({'id': len(segments), 'page': page_num, 'text': "\n\n".join(chunk)})

        full_text = "\n\n".join(s['text'] for s in segments)
        return {
            'text': full_text,
            'segments': segments,
            'metadata': {
                'paragraphs': len(paras),
                'logical_pages': len(segments),
                'method': 'docx_grouped'
            }
        }
    except Exception as e:
        return {'text': '', 'segments': [], 'metadata': {'error': str(e)}}


def route_file(file_path, ocr=False, max_units=None):
    """Route a saved file path to the appropriate extractor."""
    p = Path(file_path)
    ext = p.suffix.lower()
    result = None
    if ext == '.pdf':
        result = extract_from_pdf(p, ocr=ocr)
    elif ext == '.pptx':
        result = extract_from_pptx(p)
    elif ext == '.docx':
        result = extract_from_docx(p)
    else:
        return {'text': '', 'segments': [], 'metadata': {'error': 'unsupported file type'}}

    # Optional limit: process one page/slide/unit at a time (or first N units)
    if max_units and isinstance(max_units, int) and max_units > 0:
        segments = result.get('segments', [])[:max_units]
        result['segments'] = segments
        result['text'] = "\n\n".join(s.get('text', '') for s in segments if s.get('text'))
        result.setdefault('metadata', {})['limited_units'] = len(segments)

    return result
